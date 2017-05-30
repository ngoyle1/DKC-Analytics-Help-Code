from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.chart import XL_LABEL_POSITION


prs = Presentation("/Users/harrocyranka/Desktop/python_pptx_files/test_graph.pptx")

slide = prs.slides.add_slide(prs.slide_master.slide_layouts[0])


for shape in slide.placeholders:
    print('%d %s' % (shape.placeholder_format.idx, shape.name))

##
chart_data = ChartData()
chart_data.categories = ['West', 'East', 'North', 'South', 'Other']
chart_data.add_series('Series 1', (0.135, 0.324, 0.180, 0.235, 0.126))
placeholder = slide.placeholders[13]




chart = placeholder.insert_chart(
    XL_CHART_TYPE.PIE,chart_data
).chart

chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False

chart.plots[0].has_data_labels = True
data_labels = chart.plots[0].data_labels
data_labels.number_format = '0%'
data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

prs.save("test_graph_2.pptx")






