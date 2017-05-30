from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pandas
import numpy as np
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.chart import XL_LABEL_POSITION


##Geography Text Box Location##

###Define Textbox Location##
width = Inches(3.02)
height = Inches(0.4)
left = Inches(7.7)
top = Inches(1.6)


###
prs = Presentation('/Users/harrocyranka/Desktop/python_pptx_files/test_comparison_4.pptx')
title_slide_layout = prs.slide_layouts[2]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
title.text = raw_input("Define Title Slide: ")

##Basic Stuff
define_file_to_save = raw_input("Define File to Save: ")
define_number_of_social_media = int(raw_input("Enter Number of Social Media: "))
first_social_media = raw_input("Enter First Social Media: ")

##Geography Slide
geography_slide_layout = prs.slide_layouts[5]
geography_slide = prs.slides.add_slide(geography_slide_layout)
txBox1 = geography_slide.shapes.add_textbox(left, top, width, height)
tf = txBox1.text_frame

p = tf.add_paragraph()
p.text = first_social_media + " Proportional Audience Comparison"
p.font.size = Pt(9)
p.font.bold = True


###Create the Slides##

first_audience = raw_input("Define First Audience: ")
second_audience = raw_input("Define Second Audience: ")
first_audience_size = raw_input("Define First Audience Size: ")
second_audience_size = raw_input("Define Second Audience Size: ")
first_workbook = raw_input("Workbook with first data: ")
second_workbook = raw_input("Workbook with second data: ")
topics = ['Age', 'Income','Gender', 'Education', 'Ethnicity','Ideology']
worksheets = ["age_pivot2","income_pivot2", "gender_fixed", "education_pivot2","ethnicity", "ideology_fixed"]
slides_titles = first_audience + " and " + second_audience + " " + first_social_media + " Audience Comparison: "
n_columns = 4
nrows = [6,6,3,4,5,6]

##Enter First Social Media
for i in range(0,len(topics)):
	slide_1 = prs.slides.add_slide(prs.slide_master.slide_layouts[0])
	slide_1.placeholders[0].text = slides_titles + topics[i]
	slide_1.placeholders[10].text = first_audience
	slide_1.placeholders[11].text = second_audience
	slide_1.placeholders[12].text = "Comparison"
	slide_1.placeholders[16].text = "Total US: " + first_audience_size
	slide_1.placeholders[17].text = "Total US: " + second_audience_size


	##ADD First Graph
	data1 = pandas.read_excel(first_workbook,sheetname = worksheets[i])
	chart_data = ChartData()
	chart_data.categories = list(np.array(data1.ix[:,0:1].values.tolist()).flatten())
	chart_data.add_series(topics[i],list(np.array(data1.ix[:,2:3].values.tolist()).flatten()))
	placeholder = slide_1.placeholders[18]
	chart = placeholder.insert_chart(
	    XL_CHART_TYPE.PIE,chart_data
	).chart

	chart.has_legend = True
	chart.legend.position = XL_LEGEND_POSITION.BOTTOM
	chart.legend.include_in_layout = False
	chart.legend.font.size = Pt(8)

	chart.plots[0].has_data_labels = True
	data_labels = chart.plots[0].data_labels
	data_labels.number_format = '0.00%'
	data_labels.font.size = Pt(10)
	data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

	##ADD Second Graph
	data2 = pandas.read_excel(second_workbook,sheetname = worksheets[i])
	chart_data = ChartData()
	chart_data.categories = list(np.array(data2.ix[:,0:1].values.tolist()).flatten())
	chart_data.add_series(topics[i],list(np.array(data2.ix[:,2:3].values.tolist()).flatten()))
	placeholder = slide_1.placeholders[19]
	chart = placeholder.insert_chart(
	    XL_CHART_TYPE.PIE,chart_data
	).chart

	chart.has_legend = True
	chart.legend.position = XL_LEGEND_POSITION.BOTTOM
	chart.legend.include_in_layout = False
	chart.legend.font.size = Pt(8)

	chart.plots[0].has_data_labels = True
	data_labels = chart.plots[0].data_labels
	data_labels.number_format = '0.00%'
	data_labels.font.size = Pt(10)
	data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

	##Add Table
	table = slide_1.shapes.placeholders[20].insert_table(nrows[i],n_columns)
	table.table.cell(0, 0).text = 'Category'
	cell = table.table.rows[0].cells[0]
	paragraph = cell.text_frame.paragraphs[0]
	paragraph.font.size = Pt(10)
	paragraph = cell.text_frame.paragraphs[0]

	table.table.cell(0, 1).text = first_audience
	cell = table.table.rows[0].cells[1]
	paragraph = cell.text_frame.paragraphs[0]
	paragraph.font.size = Pt(10)
	paragraph = cell.text_frame.paragraphs[0]

	table.table.cell(0, 2).text = second_audience
	cell = table.table.rows[0].cells[2]
	paragraph = cell.text_frame.paragraphs[0]
	paragraph.font.size = Pt(10)
	paragraph = cell.text_frame.paragraphs[0]

	table.table.cell(0, 3).text = "Proportional Comparison"
	cell = table.table.rows[0].cells[3]
	paragraph = cell.text_frame.paragraphs[0]
	paragraph.font.size = Pt(10)
	paragraph = cell.text_frame.paragraphs[0]


if define_number_of_social_media == 2:
	second_social_media = raw_input("Enter Second Social Media: ")
	second_social_media_audience_1 = raw_input("Size in Second Social Media of First Group: ")
	second_social_media_audience_2 = raw_input("Size in Second Social Media of Second Group: ")
	third_workbook = raw_input("Workbook with third data: ")
	fourth_workbook = raw_input("Workbook with fourth data: ")
	new_slides_titles = first_audience + " and " + second_audience + " " + second_social_media + " Audience Comparison: "
	geography_slide = prs.slides.add_slide(geography_slide_layout)

	txBox2 = geography_slide.shapes.add_textbox(left, top, width, height)
	tf = txBox2.text_frame

	p = tf.add_paragraph()
	p.text = second_social_media + " Proportional Audience Comparison"
	p.font.size = Pt(9)
	p.font.bold = True
	for i in range(0,len(topics)):
		slide_1 = prs.slides.add_slide(prs.slide_master.slide_layouts[0])
		slide_1.placeholders[0].text = new_slides_titles + topics[i]
		slide_1.placeholders[10].text = first_audience
		slide_1.placeholders[11].text = second_audience
		slide_1.placeholders[12].text = "Comparison"
		slide_1.placeholders[16].text = "Total US: " + second_social_media_audience_1
		slide_1.placeholders[17].text = "Total US: " + second_social_media_audience_2

		##ADD First Graph
		data1 = pandas.read_excel(third_workbook,sheetname = worksheets[i])
		chart_data = ChartData()
		chart_data.categories = list(np.array(data1.ix[:,0:1].values.tolist()).flatten())
		chart_data.add_series(topics[i],list(np.array(data1.ix[:,2:3].values.tolist()).flatten()))
		placeholder = slide_1.placeholders[18]
		chart = placeholder.insert_chart(
		    XL_CHART_TYPE.PIE,chart_data
		).chart

		chart.has_legend = True
		chart.legend.position = XL_LEGEND_POSITION.BOTTOM
		chart.legend.include_in_layout = False
		chart.legend.font.size = Pt(8)

		chart.plots[0].has_data_labels = True
		data_labels = chart.plots[0].data_labels
		data_labels.font.size = Pt(10)
		data_labels.number_format = '0.00%'
		data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

		##ADD Second Graph
		data2 = pandas.read_excel(fourth_workbook,sheetname = worksheets[i])
		chart_data = ChartData()
		chart_data.categories = list(np.array(data2.ix[:,0:1].values.tolist()).flatten())
		chart_data.add_series(topics[i],list(np.array(data2.ix[:,2:3].values.tolist()).flatten()))
		placeholder = slide_1.placeholders[19]
		chart = placeholder.insert_chart(
		    XL_CHART_TYPE.PIE,chart_data
		).chart

		chart.has_legend = True
		chart.legend.position = XL_LEGEND_POSITION.BOTTOM
		chart.legend.include_in_layout = False
		chart.legend.font.size = Pt(8)

		chart.plots[0].has_data_labels = True
		data_labels = chart.plots[0].data_labels
		data_labels.number_format = '0.00%'
		data_labels.font.size = Pt(10)
		data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

		##Add Table
		table = slide_1.shapes.placeholders[20].insert_table(nrows[i],n_columns)
		table.table.cell(0, 0).text = 'Category'
		cell = table.table.rows[0].cells[0]
		paragraph = cell.text_frame.paragraphs[0]
		paragraph.font.size = Pt(10)
		paragraph = cell.text_frame.paragraphs[0]

		table.table.cell(0, 1).text = first_audience
		cell = table.table.rows[0].cells[1]
		paragraph = cell.text_frame.paragraphs[0]
		paragraph.font.size = Pt(10)
		paragraph = cell.text_frame.paragraphs[0]

		table.table.cell(0, 2).text = second_audience
		cell = table.table.rows[0].cells[2]
		paragraph = cell.text_frame.paragraphs[0]
		paragraph.font.size = Pt(10)
		paragraph = cell.text_frame.paragraphs[0]

		table.table.cell(0, 3).text = "Proportional Comparison"
		cell = table.table.rows[0].cells[3]
		paragraph = cell.text_frame.paragraphs[0]
		paragraph.font.size = Pt(10)
		paragraph = cell.text_frame.paragraphs[0]

		prs.save(define_file_to_save)

else: 
	prs.save(define_file_to_save)