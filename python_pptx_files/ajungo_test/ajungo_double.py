from pptx import Presentation
from pptx.enum.text import PP_ALIGN
import pandas

##Bring the Files with the categories
csv_to_read = raw_input("Enter file with categories: ")
x = pandas.read_csv(csv_to_read)
categories = x["x"].tolist()
categories


##

##Create the PPTX Presentation
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
insert_title = raw_input("Define Title: ")
title.text = insert_title


first_ajungo = raw_input("First Ajungo: ")
second_ajungo = raw_input("Second Ajungo: ")

###
for i in range(0,len(categories)):
    slide_1 = prs.slides.add_slide(prs.slide_layouts[4])
    slide_1.placeholders[0].text = categories[i]
    slide_1.placeholders[1].text = first_ajungo
    slide_1.placeholders[3].text = second_ajungo
    print i


file_to_save = raw_input("Write file name: ")
prs.save(file_to_save)