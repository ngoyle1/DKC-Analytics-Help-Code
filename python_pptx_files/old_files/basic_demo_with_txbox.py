from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title


number_of_social_media = int(raw_input("Number of Social Media: "))
first_social_media = raw_input("First Social Media: ")
audience = raw_input("Define Audience: ")
file_to_save = raw_input("Write file name: ")
title.text = raw_input("Define Title Slide: ")
first_social_media_audience = raw_input("First Social Media Audience: ")

###Define Textbox Location##
width = Inches(3.02)
height = Inches(0.4)
left = Inches(4.09)
top = Inches(6.6)

###Geography
slide_1 = prs.slides.add_slide(prs.slide_layouts[4])

slide_1.placeholders[0].text = audience + " " +first_social_media + " Audience: " +"Geography"
slide_1.placeholders[1].text = "Bubble Map"
slide_1.placeholders[3].text = "Tree Map"
txBox1 = slide_1.shapes.add_textbox(left, top, width, height)
tf = txBox1.text_frame

p = tf.add_paragraph()
p.text = "Total US " + first_social_media + " : " + first_social_media_audience
p.font.size = Pt(9)
p.font.bold = True



###Age and Income Slides##
slide_2 = prs.slides.add_slide(prs.slide_layouts[4])
slide_2.placeholders[0].text = audience+ " " +first_social_media + " Audience: " +"Age and Income"
slide_2.placeholders[1].text = "Age"
slide_2.placeholders[3].text = "Income"
txBox2 = slide_2.shapes.add_textbox(left, top, width, height)
tf = txBox2.text_frame

p = tf.add_paragraph()
p.text = "Total US " + first_social_media + " : " + first_social_media_audience
p.font.size = Pt(9)
p.font.bold = True



##Gender and Education
slide_3 = prs.slides.add_slide(prs.slide_layouts[4])
slide_3.placeholders[0].text = audience+ " " +first_social_media + " Audience: " + "Gender and Education"
slide_3.placeholders[1].text = "Gender"
slide_3.placeholders[3].text = "Education"
txBox3 = slide_3.shapes.add_textbox(left, top, width, height)
tf = txBox3.text_frame

p = tf.add_paragraph()
p.text = "Total US " + first_social_media + " : " + first_social_media_audience
p.font.size = Pt(9)
p.font.bold = True

##Ethnicity and Ideology
slide_4 = prs.slides.add_slide(prs.slide_layouts[4])
slide_4.placeholders[0].text = audience+ " " +first_social_media + " Audience: " + "Ethnicity and Ideology"
slide_4.placeholders[1].text = "Ethnicity"
slide_4.placeholders[3].text = "Ideology"
txBox4 = slide_4.shapes.add_textbox(left, top, width, height)
tf = txBox4.text_frame

p = tf.add_paragraph()
p.text = "Total US " + first_social_media + " : " + first_social_media_audience
p.font.size = Pt(9)
p.font.bold = True


##Check Number of Social Media
if number_of_social_media == 2:
	second_social_media = raw_input("Second Social Media: ")
	second_title = raw_input("Second Slide Title: ")
	second_social_media_audience = raw_input("Second Social Media Audience: ")

	second_title_slide = prs.slides.add_slide(prs.slide_layouts[0])
	title_2 = second_title_slide.shapes.title
	title_2.text = second_title

	slide_5 = prs.slides.add_slide(prs.slide_layouts[4])
	slide_5.placeholders[0].text = audience + " " +second_social_media + " Audience: " +"Geography"
	slide_5.placeholders[1].text = "Bubble Map"
	slide_5.placeholders[3].text = "Tree Map"
	txBox5 = slide_5.shapes.add_textbox(left, top, width, height)
	tf = txBox5.text_frame

	p = tf.add_paragraph()
	p.text = "Total US " + second_social_media + " : " + second_social_media_audience
	p.font.size = Pt(9)
	p.font.bold = True

	slide_6 = prs.slides.add_slide(prs.slide_layouts[4])
	slide_6.placeholders[0].text = audience+ " " +second_social_media + " Audience: " +"Age and Income"
	slide_6.placeholders[1].text = "Age"
	slide_6.placeholders[3].text = "Income"
	txBox6 = slide_6.shapes.add_textbox(left, top, width, height)
	tf = txBox6.text_frame

	p = tf.add_paragraph()
	p.text = "Total US " + second_social_media + " : " + second_social_media_audience
	p.font.size = Pt(9)
	p.font.bold = True

	slide_7 = prs.slides.add_slide(prs.slide_layouts[4])
	slide_7.placeholders[0].text = audience+ " " +second_social_media + " Audience: " +"Gender and Education"
	slide_7.placeholders[1].text = "Gender"
	slide_7.placeholders[3].text = "Education"
	txBox7 = slide_7.shapes.add_textbox(left, top, width, height)
	tf = txBox7.text_frame

	p = tf.add_paragraph()
	p.text = "Total US " + second_social_media + " : " + second_social_media_audience
	p.font.size = Pt(9)
	p.font.bold = True

	slide_8 = prs.slides.add_slide(prs.slide_layouts[4])
	slide_8.placeholders[0].text = audience+ " " +second_social_media + " Audience: " +"Ethnicity and Ideology"
	slide_8.placeholders[1].text = "Ethnicity"
	slide_8.placeholders[3].text = "Ideology"
	txBox8 = slide_8.shapes.add_textbox(left, top, width, height)
	tf = txBox8.text_frame

	p = tf.add_paragraph()
	p.text = "Total US " + second_social_media + " : " + second_social_media_audience
	p.font.size = Pt(9)
	p.font.bold = True

	prs.save(file_to_save)
else:
	prs.save(file_to_save)


