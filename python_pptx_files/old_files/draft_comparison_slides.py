from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

###Define Textbox Location##
width = Inches(3.02)
height = Inches(0.4)
left = Inches(4.09)
top = Inches(6.6)

###
prs = Presentation('test_comparison_4.pptx')
title_slide_layout = prs.slide_master.slide_layouts[2]
slide = prs.slides.add_slide(title_slide_layout)


for shape in slide.placeholders:
    print('%d %s' % (shape.placeholder_format.idx, shape.name))
    
slide.placeholders[0].text = "Test 1"
slide.placeholders[10].text = "Citibank"
slide.placeholders[11].text = "Citigroup"
slide.placeholders[12].text = "Comparison"
slide.placeholders[16].text = "Total US: 200,000,000"
slide.placeholders[17].text = "Total US: 200,000,000"

