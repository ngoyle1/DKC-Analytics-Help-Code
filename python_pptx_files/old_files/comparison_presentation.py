from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt



##Geography Text Box Location##

###Define Textbox Location##
width = Inches(3.02)
height = Inches(0.4)
left = Inches(7.7)
top = Inches(1.6)


###
prs = Presentation('/Users/harrocyranka/Desktop/python_pptx_files/test_comparison_4.pptx')
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
title.text = raw_input("Define Title Slide: ")

##Basic Stuff
define_file_to_save = raw_input("Define File to Save: ")
define_number_of_social_media = int(raw_input("Enter Number of Social Media: "))
first_social_media = raw_input("Enter First Social Media: ")

##Geography Slide
geography_slide_layout = prs.slide_layouts[5]
geography_slide = prs.slides.add_slide(geography_slide_layout)
txBox1 = geography_slide.shapes.add_textbox(left, top, width, height)
tf = txBox1.text_frame

p = tf.add_paragraph()
p.text = first_social_media + " Proportional Audience Comparison"
p.font.size = Pt(9)
p.font.bold = True


###Create the Slides##

first_audience = raw_input("Define First Audience: ")
second_audience = raw_input("Define Second Audience: ")
first_audience_size = raw_input("Define First Audience Size: ")
second_audience_size = raw_input("Define Second Audience Size: ")
topics = ['Age', 'Income','Gender', 'Education', 'Ethnicity','Ideology']
slides_titles = first_audience + " and " + second_audience + " " + first_social_media + " Audience Comparison: "


##Enter First Social Media
for i in range(0,len(topics)):
	slide_1 = prs.slides.add_slide(prs.slide_master.slide_layouts[0])
	slide_1.placeholders[0].text = slides_titles + topics[i]
	slide_1.placeholders[10].text = first_audience
	slide_1.placeholders[11].text = second_audience
	slide_1.placeholders[12].text = "Comparison"
	slide_1.placeholders[16].text = "Total US: " + first_audience_size
	slide_1.placeholders[17].text = "Total US: " + second_audience_size


if define_number_of_social_media == 2:
	second_social_media = raw_input("Enter Second Social Media: ")
	second_social_media_audience_1 = raw_input("Size in Second Social Media of First Group: ")
	second_social_media_audience_2 = raw_input("Size in Second Social Media of Second Group: ")
	new_slides_titles = first_audience + " and " + second_audience + " " + second_social_media + " Audience Comparison: "
	geography_slide = prs.slides.add_slide(geography_slide_layout)

	txBox2 = geography_slide.shapes.add_textbox(left, top, width, height)
	tf = txBox2.text_frame

	p = tf.add_paragraph()
	p.text = second_social_media + " Proportional Audience Comparison"
	p.font.size = Pt(9)
	p.font.bold = True
	for i in range(0,len(topics)):
		slide_1 = prs.slides.add_slide(prs.slide_master.slide_layouts[0])
		slide_1.placeholders[0].text = new_slides_titles + topics[i]
		slide_1.placeholders[10].text = first_audience
		slide_1.placeholders[11].text = second_audience
		slide_1.placeholders[12].text = "Comparison"
		slide_1.placeholders[16].text = "Total US: " + second_social_media_audience_1
		slide_1.placeholders[17].text = "Total US: " + second_social_media_audience_2

		prs.save(define_file_to_save)
else: 
	prs.save(define_file_to_save)