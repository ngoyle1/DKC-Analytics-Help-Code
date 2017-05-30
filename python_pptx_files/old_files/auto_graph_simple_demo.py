from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pandas
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.chart import XL_LABEL_POSITION

prs = Presentation("/Users/harrocyranka/Desktop/python_pptx_files/test_graph.pptx")
title_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title


first_social_media = raw_input("First Social Media: ")
audience = raw_input("Define Audience: ")
file_to_save = raw_input("Write file name: ")
title.text = raw_input("Define Title Slide: ")
first_social_media_audience = raw_input("First Social Media Audience: ")
first_social_media_data = raw_input("First Social Media Data: ")

#####Pandas Data Frame
age = pandas.read_excel(first_social_media_data,sheetname = "age_pivot2")
income = pandas.read_excel(first_social_media_data,sheetname = "income_pivot2")
gender = pandas.read_excel(first_social_media_data,sheetname = "gender_fixed")
education = pandas.read_excel(first_social_media_data,sheetname = "education_pivot2")
ethnicity = pandas.read_excel(first_social_media_data,sheetname = "ethnicity")
ideology = pandas.read_excel(first_social_media_data,sheetname = "ideology_fixed")


###Define Textbox Location##
width = Inches(3.02)
height = Inches(0.4)
left = Inches(4.09)
top = Inches(6.6)

###Geography
slide_1 = prs.slides.add_slide(prs.slide_master.slide_layouts[0])

slide_1.placeholders[0].text = audience + " " +first_social_media + " Audience: " +"Geography"
slide_1.placeholders[1].text = "Bubble Map"
slide_1.placeholders[3].text = "Tree Map"
txBox1 = slide_1.shapes.add_textbox(left, top, width, height)
tf = txBox1.text_frame

p = tf.add_paragraph()
p.text = "Total US " + first_social_media + " : " + first_social_media_audience
p.font.size = Pt(9)
p.font.bold = True



###Age and Income Slides##
slide_2 = prs.slides.add_slide(prs.slide_master.slide_layouts[0])
slide_2.placeholders[0].text = audience+ " " +first_social_media + " Audience: " +"Age and Income"
slide_2.placeholders[1].text = "Age"
slide_2.placeholders[3].text = "Income"
txBox2 = slide_2.shapes.add_textbox(left, top, width, height)
tf = txBox2.text_frame

p = tf.add_paragraph()
p.text = "Total US " + first_social_media + " : " + first_social_media_audience
p.font.size = Pt(9)
p.font.bold = True


##Add Age Chart
chart_data = ChartData()
chart_data.categories = age["age_categories"].tolist()
chart_data.add_series("Age",age["proportional"].tolist())
placeholder = slide_2.placeholders[13]


chart = placeholder.insert_chart(
    XL_CHART_TYPE.PIE,chart_data
).chart

chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False
chart.legend.font.size = Pt(8)

chart.plots[0].has_data_labels = True
data_labels = chart.plots[0].data_labels
data_labels.number_format = '0.00%'
data_labels.font.size = Pt(10)
data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

##Add Income Chart
chart_data = ChartData()
chart_data.categories = income["income_categories"].tolist()
chart_data.add_series("Income",income["proportional"].tolist())
placeholder = slide_2.placeholders[14]


chart = placeholder.insert_chart(
    XL_CHART_TYPE.PIE,chart_data
).chart

chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False
chart.legend.font.size = Pt(8)

chart.plots[0].has_data_labels = True
data_labels = chart.plots[0].data_labels
data_labels.font.size = Pt(10)
data_labels.number_format = '0.00%'
data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

##Gender and Education
slide_3 = prs.slides.add_slide(prs.slide_master.slide_layouts[0])
slide_3.placeholders[0].text = audience+ " " +first_social_media + " Audience: " + "Gender and Education"
slide_3.placeholders[1].text = "Gender"
slide_3.placeholders[3].text = "Education"
txBox3 = slide_3.shapes.add_textbox(left, top, width, height)
tf = txBox3.text_frame

p = tf.add_paragraph()
p.text = "Total US " + first_social_media + " : " + first_social_media_audience
p.font.size = Pt(9)
p.font.bold = True


##Add Gender Chart
chart_data = ChartData()
chart_data.categories = gender["Gender"].tolist()
chart_data.add_series("Gender",gender['percent'].tolist())
placeholder = slide_3.placeholders[13]


chart = placeholder.insert_chart(
    XL_CHART_TYPE.PIE,chart_data
).chart

chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False
chart.legend.font.size = Pt(8)

chart.plots[0].has_data_labels = True
data_labels = chart.plots[0].data_labels
data_labels.number_format = '0.00%'
data_labels.font.size = Pt(10)
data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

##Add Education Chart
chart_data = ChartData()
chart_data.categories = education["categories"].tolist()
chart_data.add_series("Education",education['proportional'].tolist())
placeholder = slide_3.placeholders[14]


chart = placeholder.insert_chart(
    XL_CHART_TYPE.PIE,chart_data
).chart

chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False
chart.legend.font.size = Pt(8)

chart.plots[0].has_data_labels = True
data_labels = chart.plots[0].data_labels
data_labels.number_format = '0.00%'
data_labels.font.size = Pt(10)
data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

##Ethnicity and Ideology
slide_4 = prs.slides.add_slide(prs.slide_master.slide_layouts[0])
slide_4.placeholders[0].text = audience+ " " +first_social_media + " Audience: " + "Ethnicity and Ideology"
slide_4.placeholders[1].text = "Ethnicity"
slide_4.placeholders[3].text = "Ideology"
txBox4 = slide_4.shapes.add_textbox(left, top, width, height)
tf = txBox4.text_frame

p = tf.add_paragraph()
p.text = "Total US " + first_social_media + " : " + first_social_media_audience
p.font.size = Pt(9)
p.font.bold = True

##Add Ethnicity Chart
chart_data = ChartData()
chart_data.categories = ethnicity["Ethnic/Cultural Group"].tolist()
chart_data.add_series("Ethnicity",ethnicity['percent'].tolist())
placeholder = slide_4.placeholders[13]


chart = placeholder.insert_chart(
    XL_CHART_TYPE.PIE,chart_data
).chart

chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False
chart.legend.font.size = Pt(8)

chart.plots[0].has_data_labels = True
data_labels = chart.plots[0].data_labels
data_labels.number_format = '0.00%'
data_labels.font.size = Pt(10)
data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

##Add Ideology Graph
chart_data = ChartData()
chart_data.categories = ideology["Politics"].tolist()
chart_data.add_series("Ideology",ideology['proportional'].tolist())
placeholder = slide_4.placeholders[14]


chart = placeholder.insert_chart(
    XL_CHART_TYPE.PIE,chart_data
).chart

chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False
chart.legend.font.size = Pt(8)

chart.plots[0].has_data_labels = True
data_labels = chart.plots[0].data_labels
data_labels.number_format = '0.00%'
data_labels.font.size = Pt(10)
data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

prs.save(file_to_save)



